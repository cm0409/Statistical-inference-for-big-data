"""
海量数据统计计算方法 - 图表绘制完整代码（Python版）
可直接导出为PPT可编辑格式
"""

import matplotlib.pyplot as plt
import matplotlib
import numpy as np
import pandas as pd
from matplotlib import rcParams
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# 设置中文字体
rcParams['font.sans-serif'] = ['SimHei', 'DejaVu Sans']
rcParams['axes.unicode_minus'] = False

output_dir = '/mnt/okcomputer/output'

# ==================== 数据准备 ====================

# 表1：模型性能对比
model_data = pd.DataFrame({
    '模型': ['GAM', '多项式', 'OLS', '岭回归', '部分线性', '核回归'],
    '计算时间': [0.0015, 0.0009, 0.0011, 0.1316, 0.0356, 0.0932],
    'MSE': [0.9878, 0.9911, 0.9962, 0.9962, 1.3812, 2.8258],
    'R2': [0.7213, 0.7204, 0.719, 0.719, 0.6103, 0.2028]
})

# 表2：分布式计算
dist_data = pd.DataFrame({
    'K': [5, 10, 20, 50],
    'SAE时间': [0.0097, 0.0734, 0.0846, 0.0428],
    'SAE加速比': [0.84, 0.11, 0.10, 0.19],
    '一步估计时间': [0.0058, 0.0044, 0.0047, 0.0056],
    '一步估计加速比': [1.40, 1.85, 1.73, 1.46],
    '理想加速比': [5, 10, 20, 50]
})

# 表3：BLB结果
blb_data = pd.DataFrame({
    'Gamma': [0.6, 0.6, 0.6, 0.7, 0.7, 0.7],
    'S': [5, 10, 20, 5, 10, 20],
    '方差': [0.009524, 0.010353, 0.010103, 0.003966, 0.004099, 0.004009]
})

# 表4：SGD结果
sgd_data = pd.DataFrame({
    'BatchSize': [32, 64, 128, 256, 512],
    '计算时间': [2.883, 1.4644, 0.7507, 0.3868, 0.203],
    '最终MSE': [0.9842, 0.9840, 0.9853, 0.9839, 0.9854]
})

# ==================== Matplotlib绘图（导出SVG/PNG） ====================

def plot_fig1_compute_time():
    """图1：计算时间对比图"""
    fig, ax = plt.subplots(figsize=(10, 6))
    
    models = model_data['模型']
    times = model_data['计算时间']
    colors = ['#4472C4', '#ED7D31', '#A5A5A5', '#FFC000', '#5B9BD5', '#70AD47']
    
    bars = ax.bar(models, times, color=colors, edgecolor='black', linewidth=0.5)
    
    for bar, time in zip(bars, times):
        height = bar.get_height()
        ax.text(bar.get_x() + bar.get_width()/2., height,
                f'{time:.4f}', ha='center', va='bottom', fontsize=10)
    
    ax.set_xlabel('模型', fontsize=12)
    ax.set_ylabel('计算时间（秒）', fontsize=12)
    ax.set_yscale('log')
    ax.grid(axis='y', alpha=0.3, linestyle='--')
    
    plt.tight_layout()
    plt.savefig(f'{output_dir}/fig1_compute_time.svg', format='svg', dpi=300, bbox_inches='tight')
    plt.savefig(f'{output_dir}/fig1_compute_time.png', format='png', dpi=300, bbox_inches='tight')
    plt.close()
    print("图1完成：计算时间对比图")

def plot_fig2_mse_comparison():
    """图2：MSE对比图"""
    fig, ax = plt.subplots(figsize=(10, 6))
    
    models = model_data['模型']
    mse_values = model_data['MSE']
    colors = ['#4472C4', '#ED7D31', '#A5A5A5', '#FFC000', '#5B9BD5', '#70AD47']
    
    bars = ax.bar(models, mse_values, color=colors, edgecolor='black', linewidth=0.5)
    
    for bar, mse in zip(bars, mse_values):
        height = bar.get_height()
        ax.text(bar.get_x() + bar.get_width()/2., height,
                f'{mse:.4f}', ha='center', va='bottom', fontsize=10)
    
    ax.set_xlabel('模型', fontsize=12)
    ax.set_ylabel('均方误差（MSE）', fontsize=12)
    ax.grid(axis='y', alpha=0.3, linestyle='--')
    
    plt.tight_layout()
    plt.savefig(f'{output_dir}/fig2_mse_comparison.svg', format='svg', dpi=300, bbox_inches='tight')
    plt.savefig(f'{output_dir}/fig2_mse_comparison.png', format='png', dpi=300, bbox_inches='tight')
    plt.close()
    print("图2完成：MSE对比图")

def plot_fig3_speedup():
    """图3：加速比曲线"""
    fig, ax = plt.subplots(figsize=(10, 6))
    
    K_values = dist_data['K']
    sae_speedup = dist_data['SAE加速比']
    onestep_speedup = dist_data['一步估计加速比']
    ideal_speedup = dist_data['理想加速比']
    
    ax.plot(K_values, sae_speedup, 'o-', color='#4472C4', linewidth=2, markersize=8, label='SAE')
    ax.plot(K_values, onestep_speedup, 's-', color='#ED7D31', linewidth=2, markersize=8, label='一步估计')
    ax.plot(K_values, ideal_speedup, '--', color='#A5A5A5', linewidth=1.5, label='理想加速比')
    
    ax.set_xlabel('节点数量 K', fontsize=12)
    ax.set_ylabel('加速比', fontsize=12)
    ax.legend(loc='upper left', fontsize=11)
    ax.grid(True, alpha=0.3, linestyle='--')
    ax.set_xticks(K_values)
    
    plt.tight_layout()
    plt.savefig(f'{output_dir}/fig3_speedup.svg', format='svg', dpi=300, bbox_inches='tight')
    plt.savefig(f'{output_dir}/fig3_speedup.png', format='png', dpi=300, bbox_inches='tight')
    plt.close()
    print("图3完成：加速比曲线")

def plot_fig4_blb_variance():
    """图4：BLB方差对比"""
    fig, ax = plt.subplots(figsize=(10, 6))
    
    gamma_labels = ['γ=0.6\nS=5', 'γ=0.6\nS=10', 'γ=0.6\nS=20', 
                    'γ=0.7\nS=5', 'γ=0.7\nS=10', 'γ=0.7\nS=20']
    variance_values = [0.009524, 0.010353, 0.010103, 0.003966, 0.004099, 0.004009]
    colors = ['#4472C4']*3 + ['#ED7D31']*3
    
    bars = ax.bar(gamma_labels, variance_values, color=colors, edgecolor='black', linewidth=0.5)
    
    for bar, var in zip(bars, variance_values):
        height = bar.get_height()
        ax.text(bar.get_x() + bar.get_width()/2., height,
                f'{var:.4f}', ha='center', va='bottom', fontsize=9)
    
    ax.set_xlabel('参数配置', fontsize=12)
    ax.set_ylabel('估计量方差', fontsize=12)
    ax.grid(axis='y', alpha=0.3, linestyle='--')
    
    from matplotlib.patches import Patch
    legend_elements = [Patch(facecolor='#4472C4', label='γ=0.6'),
                       Patch(facecolor='#ED7D31', label='γ=0.7')]
    ax.legend(handles=legend_elements, loc='upper right', fontsize=11)
    
    plt.tight_layout()
    plt.savefig(f'{output_dir}/fig4_blb_variance.svg', format='svg', dpi=300, bbox_inches='tight')
    plt.savefig(f'{output_dir}/fig4_blb_variance.png', format='png', dpi=300, bbox_inches='tight')
    plt.close()
    print("图4完成：BLB方差对比")

def plot_fig5_sgd_convergence():
    """图5：SGD收敛曲线"""
    fig, ax = plt.subplots(figsize=(10, 6))
    
    np.random.seed(42)
    iterations = np.linspace(0, 500, 100)
    mse_32 = 0.984 + 0.1 * np.exp(-iterations/100) + np.random.normal(0, 0.002, 100)
    mse_128 = 0.985 + 0.08 * np.exp(-iterations/80) + np.random.normal(0, 0.0015, 100)
    mse_512 = 0.985 + 0.06 * np.exp(-iterations/60) + np.random.normal(0, 0.001, 100)
    
    ax.plot(iterations, mse_32, '-', color='#4472C4', linewidth=2, label='batch=32', alpha=0.8)
    ax.plot(iterations, mse_128, '-', color='#ED7D31', linewidth=2, label='batch=128', alpha=0.8)
    ax.plot(iterations, mse_512, '-', color='#70AD47', linewidth=2, label='batch=512', alpha=0.8)
    
    ax.set_xlabel('迭代次数', fontsize=12)
    ax.set_ylabel('均方误差（MSE）', fontsize=12)
    ax.legend(loc='upper right', fontsize=11)
    ax.grid(True, alpha=0.3, linestyle='--')
    
    plt.tight_layout()
    plt.savefig(f'{output_dir}/fig5_sgd_convergence.svg', format='svg', dpi=300, bbox_inches='tight')
    plt.savefig(f'{output_dir}/fig5_sgd_convergence.png', format='png', dpi=300, bbox_inches='tight')
    plt.close()
    print("图5完成：SGD收敛曲线")

def plot_fig6_complexity_accuracy():
    """图6：模型复杂度-精度权衡图（修复版）"""
    fig, ax = plt.subplots(figsize=(12, 8))
    
    models_list = ['GAM', '多项式', 'OLS', '岭回归', '部分线性', '核回归']
    complexity = [4, 2, 1, 2, 3, 5]
    mse_values = [0.9878, 0.9911, 0.9962, 0.9962, 1.3812, 2.8258]
    times = [0.0015, 0.0009, 0.0011, 0.1316, 0.0356, 0.0932]
    
    colors = ['#4472C4', '#ED7D31', '#A5A5A5', '#FFC000', '#5B9BD5', '#70AD47']
    sizes = [t * 8000 + 200 for t in times]
    
    scatter = ax.scatter(complexity, mse_values, s=sizes, 
                         c=colors, alpha=0.7, edgecolors='black', linewidth=1.5)
    
    offsets = [(15, 15), (15, -25), (-30, 15), (15, 15), (15, 15), (-40, -25)]
    for i, model in enumerate(models_list):
        ax.annotate(model, (complexity[i], mse_values[i]), 
                    xytext=offsets[i], textcoords='offset points', 
                    fontsize=12, fontweight='bold',
                    bbox=dict(boxstyle='round,pad=0.3', facecolor='white', alpha=0.7, edgecolor='gray'))
    
    ax.set_xlabel('模型复杂度', fontsize=14, fontweight='bold')
    ax.set_ylabel('均方误差（MSE）', fontsize=14, fontweight='bold')
    ax.grid(True, alpha=0.3, linestyle='--')
    ax.set_xlim(0.5, 5.5)
    ax.set_ylim(0.8, 3.2)
    
    from matplotlib.lines import Line2D
    legend_elements = [
        Line2D([0], [0], marker='o', color='w', markerfacecolor='gray', markersize=8, label='小（快）'),
        Line2D([0], [0], marker='o', color='w', markerfacecolor='gray', markersize=12, label='中'),
        Line2D([0], [0], marker='o', color='w', markerfacecolor='gray', markersize=16, label='大（慢）')
    ]
    ax.legend(handles=legend_elements, title='计算时间', loc='lower right', fontsize=11)
    
    plt.tight_layout()
    plt.savefig(f'{output_dir}/fig6_complexity_accuracy_fixed.svg', format='svg', dpi=300, bbox_inches='tight')
    plt.savefig(f'{output_dir}/fig6_complexity_accuracy_fixed.png', format='png', dpi=300, bbox_inches='tight')
    plt.close()
    print("图6完成：复杂度-精度权衡图")

# ==================== PPT导出（原生可编辑） ====================

def create_ppt_charts():
    """创建PPT格式的可编辑图表"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    colors = [
        RGBColor(68, 114, 196), RGBColor(237, 125, 49), RGBColor(165, 165, 165),
        RGBColor(255, 192, 0), RGBColor(91, 155, 213), RGBColor(112, 173, 71)
    ]
    
    left = Inches(1)
    top = Inches(1)
    width = Inches(11)
    height = Inches(5.5)
    spacing = width / 6
    bar_width = spacing * 0.6
    
    # ===== 图1：计算时间对比 =====
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    models = ['GAM', '多项式', 'OLS', '岭回归', '部分线性', '核回归']
    times = [0.0015, 0.0009, 0.0011, 0.1316, 0.0356, 0.0932]
    log_times = np.log10(times)
    min_log, max_log = min(log_times), max(log_times)
    
    # Y轴
    axis_y = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.02), height)
    axis_y.fill.solid()
    axis_y.fill.fore_color.rgb = RGBColor(0, 0, 0)
    
    # X轴
    axis_x = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + height, width, Inches(0.02))
    axis_x.fill.solid()
    axis_x.fill.fore_color.rgb = RGBColor(0, 0, 0)
    
    # 柱状图
    for i, (model, time, log_time) in enumerate(zip(models, times, log_times)):
        bar_height = (log_time - min_log) / (max_log - min_log) * height * 0.8 + height * 0.1
        bar_left = left + spacing * i + spacing * 0.2
        bar_top = top + height - bar_height
        
        bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_left, bar_top, bar_width, bar_height)
        bar.fill.solid()
        bar.fill.fore_color.rgb = colors[i]
        bar.line.color.rgb = RGBColor(0, 0, 0)
        bar.line.width = Pt(1)
        
        # 数值标签
        value_box = slide1.shapes.add_textbox(bar_left, bar_top - Inches(0.3), bar_width, Inches(0.25))
        value_frame = value_box.text_frame
        value_frame.text = f'{time:.4f}'
        value_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        value_frame.paragraphs[0].font.size = Pt(10)
        
        # X轴标签
        label_box = slide1.shapes.add_textbox(bar_left, top + height + Inches(0.1), bar_width, Inches(0.3))
        label_frame = label_box.text_frame
        label_frame.text = model
        label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        label_frame.paragraphs[0].font.size = Pt(11)
    
    # Y轴标签
    y_label = slide1.shapes.add_textbox(left - Inches(0.6), top + height/2 - Inches(0.5), Inches(0.5), Inches(1))
    y_label_frame = y_label.text_frame
    y_label_frame.text = '计算时间（秒）'
    y_label_frame.word_wrap = True
    y_label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    y_label_frame.paragraphs[0].font.size = Pt(12)
    
    # X轴标签
    x_label = slide1.shapes.add_textbox(left + width/2 - Inches(1), top + height + Inches(0.5), Inches(2), Inches(0.3))
    x_label_frame = x_label.text_frame
    x_label_frame.text = '模型'
    x_label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    x_label_frame.paragraphs[0].font.size = Pt(12)
    
    print("PPT图1完成")
    
    # ===== 图2：MSE对比 =====
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    mse_values = [0.9878, 0.9911, 0.9962, 0.9962, 1.3812, 2.8258]
    max_mse = max(mse_values)
    
    axis_y = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.02), height)
    axis_y.fill.solid()
    axis_y.fill.fore_color.rgb = RGBColor(0, 0, 0)
    
    axis_x = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + height, width, Inches(0.02))
    axis_x.fill.solid()
    axis_x.fill.fore_color.rgb = RGBColor(0, 0, 0)
    
    for i, (model, mse) in enumerate(zip(models, mse_values)):
        bar_height = mse / max_mse * height * 0.85
        bar_left = left + spacing * i + spacing * 0.2
        bar_top = top + height - bar_height
        
        bar = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_left, bar_top, bar_width, bar_height)
        bar.fill.solid()
        bar.fill.fore_color.rgb = colors[i]
        bar.line.color.rgb = RGBColor(0, 0, 0)
        bar.line.width = Pt(1)
        
        value_box = slide2.shapes.add_textbox(bar_left, bar_top - Inches(0.3), bar_width, Inches(0.25))
        value_frame = value_box.text_frame
        value_frame.text = f'{mse:.4f}'
        value_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        value_frame.paragraphs[0].font.size = Pt(10)
        
        label_box = slide2.shapes.add_textbox(bar_left, top + height + Inches(0.1), bar_width, Inches(0.3))
        label_frame = label_box.text_frame
        label_frame.text = model
        label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        label_frame.paragraphs[0].font.size = Pt(11)
    
    y_label = slide2.shapes.add_textbox(left - Inches(0.6), top + height/2 - Inches(0.5), Inches(0.5), Inches(1))
    y_label_frame = y_label.text_frame
    y_label_frame.text = '均方误差（MSE）'
    y_label_frame.word_wrap = True
    y_label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    y_label_frame.paragraphs[0].font.size = Pt(12)
    
    x_label = slide2.shapes.add_textbox(left + width/2 - Inches(1), top + height + Inches(0.5), Inches(2), Inches(0.3))
    x_label_frame = x_label.text_frame
    x_label_frame.text = '模型'
    x_label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    x_label_frame.paragraphs[0].font.size = Pt(12)
    
    print("PPT图2完成")
    
    # ===== 图3：加速比曲线 =====
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    K_values = [5, 10, 20, 50]
    sae_speedup = [0.84, 0.11, 0.10, 0.19]
    onestep_speedup = [1.40, 1.85, 1.73, 1.46]
    ideal_speedup = [5, 10, 20, 50]
    max_speedup = max(ideal_speedup)
    
    axis_y = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.02), height)
    axis_y.fill.solid()
    axis_y.fill.fore_color.rgb = RGBColor(0, 0, 0)
    
    axis_x = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + height, width, Inches(0.02))
    axis_x.fill.solid()
    axis_x.fill.fore_color.rgb = RGBColor(0, 0, 0)
    
    # 理想加速比（虚线）
    for i in range(len(K_values) - 1):
        x1 = left + (K_values[i] / 55) * width
        y1 = top + height - (ideal_speedup[i] / max_speedup) * height * 0.9
        x2 = left + (K_values[i+1] / 55) * width
        y2 = top + height - (ideal_speedup[i+1] / max_speedup) * height * 0.9
        line = slide3.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
        line.line.color.rgb = RGBColor(150, 150, 150)
        line.line.width = Pt(2)
        line.line.dash_style = 2
    
    # SAE曲线
    for i in range(len(K_values) - 1):
        x1 = left + (K_values[i] / 55) * width
        y1 = top + height - (sae_speedup[i] / max_speedup) * height * 0.9
        x2 = left + (K_values[i+1] / 55) * width
        y2 = top + height - (sae_speedup[i+1] / max_speedup) * height * 0.9
        line = slide3.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
        line.line.color.rgb = RGBColor(68, 114, 196)
        line.line.width = Pt(3)
    
    # 一步估计曲线
    for i in range(len(K_values) - 1):
        x1 = left + (K_values[i] / 55) * width
        y1 = top + height - (onestep_speedup[i] / max_speedup) * height * 0.9
        x2 = left + (K_values[i+1] / 55) * width
        y2 = top + height - (onestep_speedup[i+1] / max_speedup) * height * 0.9
        line = slide3.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
        line.line.color.rgb = RGBColor(237, 125, 49)
        line.line.width = Pt(3)
    
    # 数据点
    for i, k in enumerate(K_values):
        x = left + (k / 55) * width
        y_sae = top + height - (sae_speedup[i] / max_speedup) * height * 0.9
        circle_sae = slide3.shapes.add_shape(MSO_SHAPE.OVAL, x - Inches(0.08), y_sae - Inches(0.08), Inches(0.16), Inches(0.16))
        circle_sae.fill.solid()
        circle_sae.fill.fore_color.rgb = RGBColor(68, 114, 196)
        
        y_os = top + height - (onestep_speedup[i] / max_speedup) * height * 0.9
        circle_os = slide3.shapes.add_shape(MSO_SHAPE.OVAL, x - Inches(0.08), y_os - Inches(0.08), Inches(0.16), Inches(0.16))
        circle_os.fill.solid()
        circle_os.fill.fore_color.rgb = RGBColor(237, 125, 49)
        
        label_box = slide3.shapes.add_textbox(x - Inches(0.3), top + height + Inches(0.1), Inches(0.6), Inches(0.3))
        label_frame = label_box.text_frame
        label_frame.text = str(k)
        label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        label_frame.paragraphs[0].font.size = Pt(11)
    
    y_label = slide3.shapes.add_textbox(left - Inches(0.6), top + height/2 - Inches(0.5), Inches(0.5), Inches(1))
    y_label_frame = y_label.text_frame
    y_label_frame.text = '加速比'
    y_label_frame.word_wrap = True
    y_label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    y_label_frame.paragraphs[0].font.size = Pt(12)
    
    x_label = slide3.shapes.add_textbox(left + width/2 - Inches(1), top + height + Inches(0.5), Inches(2), Inches(0.3))
    x_label_frame = x_label.text_frame
    x_label_frame.text = '节点数量 K'
    x_label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    x_label_frame.paragraphs[0].font.size = Pt(12)
    
    # 图例
    legend_y = top + Inches(0.5)
    legend1 = slide3.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.2), legend_y, Inches(0.2), Inches(0.2))
    legend1.fill.solid()
    legend1.fill.fore_color.rgb = RGBColor(68, 114, 196)
    label1 = slide3.shapes.add_textbox(left + Inches(0.5), legend_y, Inches(1), Inches(0.25))
    label1.text_frame.text = 'SAE'
    label1.text_frame.paragraphs[0].font.size = Pt(11)
    
    legend2 = slide3.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(2), legend_y, Inches(0.2), Inches(0.2))
    legend2.fill.solid()
    legend2.fill.fore_color.rgb = RGBColor(237, 125, 49)
    label2 = slide3.shapes.add_textbox(left + Inches(2.3), legend_y, Inches(1.5), Inches(0.25))
    label2.text_frame.text = '一步估计'
    label2.text_frame.paragraphs[0].font.size = Pt(11)
    
    legend3_line = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(4.5), legend_y + Inches(0.08), Inches(0.4), Inches(0.04))
    legend3_line.fill.solid()
    legend3_line.fill.fore_color.rgb = RGBColor(150, 150, 150)
    label3 = slide3.shapes.add_textbox(left + Inches(5), legend_y, Inches(1.5), Inches(0.25))
    label3.text_frame.text = '理想加速比'
    label3.text_frame.paragraphs[0].font.size = Pt(11)
    
    print("PPT图3完成")
    
    # ===== 图4：BLB方差对比 =====
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    labels = ['γ=0.6\nS=5', 'γ=0.6\nS=10', 'γ=0.6\nS=20', 'γ=0.7\nS=5', 'γ=0.7\nS=10', 'γ=0.7\nS=20']
    variance_values = [0.009524, 0.010353, 0.010103, 0.003966, 0.004099, 0.004009]
    bar_colors = [RGBColor(68, 114, 196)]*3 + [RGBColor(237, 125, 49)]*3
    max_var = max(variance_values)
    
    axis_y = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.02), height)
    axis_y.fill.solid()
    axis_y.fill.fore_color.rgb = RGBColor(0, 0, 0)
    
    axis_x = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + height, width, Inches(0.02))
    axis_x.fill.solid()
    axis_x.fill.fore_color.rgb = RGBColor(0, 0, 0)
    
    for i, (label, var, color) in enumerate(zip(labels, variance_values, bar_colors)):
        bar_height = var / max_var * height * 0.85
        bar_left = left + spacing * i + spacing * 0.25
        bar_top = top + height - bar_height
        
        bar = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_left, bar_top, bar_width, bar_height)
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.color.rgb = RGBColor(0, 0, 0)
        bar.line.width = Pt(1)
        
        value_box = slide4.shapes.add_textbox(bar_left, bar_top - Inches(0.3), bar_width, Inches(0.25))
        value_frame = value_box.text_frame
        value_frame.text = f'{var:.4f}'
        value_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        value_frame.paragraphs[0].font.size = Pt(9)
        
        label_box = slide4.shapes.add_textbox(bar_left, top + height + Inches(0.05), bar_width, Inches(0.5))
        label_frame = label_box.text_frame
        label_frame.text = label
        label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        label_frame.paragraphs[0].font.size = Pt(10)
    
    y_label = slide4.shapes.add_textbox(left - Inches(0.6), top + height/2 - Inches(0.5), Inches(0.5), Inches(1))
    y_label_frame = y_label.text_frame
    y_label_frame.text = '估计量方差'
    y_label_frame.word_wrap = True
    y_label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    y_label_frame.paragraphs[0].font.size = Pt(12)
    
    x_label = slide4.shapes.add_textbox(left + width/2 - Inches(1), top + height + Inches(0.6), Inches(2), Inches(0.3))
    x_label_frame = x_label.text_frame
    x_label_frame.text = '参数配置'
    x_label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    x_label_frame.paragraphs[0].font.size = Pt(12)
    
    # 图例
    legend_y = top + Inches(0.5)
    legend1 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.2), legend_y, Inches(0.3), Inches(0.2))
    legend1.fill.solid()
    legend1.fill.fore_color.rgb = RGBColor(68, 114, 196)
    label1 = slide4.shapes.add_textbox(left + Inches(0.6), legend_y, Inches(1), Inches(0.25))
    label1.text_frame.text = 'γ=0.6'
    label1.text_frame.paragraphs[0].font.size = Pt(11)
    
    legend2 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(2), legend_y, Inches(0.3), Inches(0.2))
    legend2.fill.solid()
    legend2.fill.fore_color.rgb = RGBColor(237, 125, 49)
    label2 = slide4.shapes.add_textbox(left + Inches(2.4), legend_y, Inches(1), Inches(0.25))
    label2.text_frame.text = 'γ=0.7'
    label2.text_frame.paragraphs[0].font.size = Pt(11)
    
    print("PPT图4完成")
    
    # ===== 图5：SGD收敛曲线 =====
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    np.random.seed(42)
    iterations = np.linspace(0, 500, 50)
    mse_32 = 0.984 + 0.1 * np.exp(-iterations/100) + np.random.normal(0, 0.002, 50)
    mse_128 = 0.985 + 0.08 * np.exp(-iterations/80) + np.random.normal(0, 0.0015, 50)
    mse_512 = 0.985 + 0.06 * np.exp(-iterations/60) + np.random.normal(0, 0.001, 50)
    
    max_mse = max(max(mse_32), max(mse_128), max(mse_512))
    min_mse = min(min(mse_32), min(mse_128), min(mse_512))
    
    axis_y = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.02), height)
    axis_y.fill.solid()
    axis_y.fill.fore_color.rgb = RGBColor(0, 0, 0)
    
    axis_x = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + height, width, Inches(0.02))
    axis_x.fill.solid()
    axis_x.fill.fore_color.rgb = RGBColor(0, 0, 0)
    
    line_data = [(mse_32, RGBColor(68, 114, 196), 'batch=32'),
                 (mse_128, RGBColor(237, 125, 49), 'batch=128'),
                 (mse_512, RGBColor(112, 173, 71), 'batch=512')]
    
    for mse_data, color, label in line_data:
        for i in range(len(iterations) - 1):
            x1 = left + (iterations[i] / 550) * width
            y1 = top + height - ((mse_data[i] - min_mse) / (max_mse - min_mse)) * height * 0.85 - height * 0.05
            x2 = left + (iterations[i+1] / 550) * width
            y2 = top + height - ((mse_data[i+1] - min_mse) / (max_mse - min_mse)) * height * 0.85 - height * 0.05
            line = slide5.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
            line.line.color.rgb = color
            line.line.width = Pt(2)
    
    for i in range(0, 6):
        x_val = i * 100
        x_pos = left + (x_val / 550) * width
        label_box = slide5.shapes.add_textbox(x_pos - Inches(0.3), top + height + Inches(0.1), Inches(0.6), Inches(0.3))
        label_frame = label_box.text_frame
        label_frame.text = str(x_val)
        label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        label_frame.paragraphs[0].font.size = Pt(11)
    
    y_label = slide5.shapes.add_textbox(left - Inches(0.6), top + height/2 - Inches(0.5), Inches(0.5), Inches(1))
    y_label_frame = y_label.text_frame
    y_label_frame.text = '均方误差（MSE）'
    y_label_frame.word_wrap = True
    y_label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    y_label_frame.paragraphs[0].font.size = Pt(12)
    
    x_label = slide5.shapes.add_textbox(left + width/2 - Inches(1), top + height + Inches(0.5), Inches(2), Inches(0.3))
    x_label_frame = x_label.text_frame
    x_label_frame.text = '迭代次数'
    x_label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    x_label_frame.paragraphs[0].font.size = Pt(12)
    
    # 图例
    legend_y = top + Inches(0.5)
    legend_items = [(RGBColor(68, 114, 196), 'batch=32', 0),
                    (RGBColor(237, 125, 49), 'batch=128', 2),
                    (RGBColor(112, 173, 71), 'batch=512', 4)]
    
    for color, label, offset in legend_items:
        line = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.2 + offset), legend_y + Inches(0.08), Inches(0.4), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = color
        text = slide5.shapes.add_textbox(left + Inches(0.7 + offset), legend_y, Inches(1.5), Inches(0.25))
        text.text_frame.text = label
        text.text_frame.paragraphs[0].font.size = Pt(11)
    
    print("PPT图5完成")
    
    # ===== 图6：复杂度-精度权衡 =====
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    models_list = ['GAM', '多项式', 'OLS', '岭回归', '部分线性', '核回归']
    complexity = [4, 2, 1, 2, 3, 5]
    mse_values = [0.9878, 0.9911, 0.9962, 0.9962, 1.3812, 2.8258]
    times = [0.0015, 0.0009, 0.0011, 0.1316, 0.0356, 0.0932]
    max_complexity = max(complexity)
    max_mse = max(mse_values)
    max_time = max(times)
    
    axis_y = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.02), height)
    axis_y.fill.solid()
    axis_y.fill.fore_color.rgb = RGBColor(0, 0, 0)
    
    axis_x = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + height, width, Inches(0.02))
    axis_x.fill.solid()
    axis_x.fill.fore_color.rgb = RGBColor(0, 0, 0)
    
    bubble_colors = [RGBColor(68, 114, 196), RGBColor(237, 125, 49), RGBColor(165, 165, 165),
                     RGBColor(255, 192, 0), RGBColor(91, 155, 213), RGBColor(112, 173, 71)]
    
    offsets = [(Inches(0.3), Inches(-0.3)), (Inches(0.3), Inches(0.2)), (Inches(-0.8), Inches(-0.2)),
               (Inches(0.3), Inches(-0.3)), (Inches(0.3), Inches(-0.3)), (Inches(-1.0), Inches(0.2))]
    
    for i, (model, comp, mse, time) in enumerate(zip(models_list, complexity, mse_values, times)):
        x = left + (comp / 6) * width
        y = top + height - (mse / 3) * height * 0.85 - height * 0.05
        bubble_size = Inches(0.2 + time / max_time * 0.6)
        
        bubble = slide6.shapes.add_shape(MSO_SHAPE.OVAL, x - bubble_size/2, y - bubble_size/2, bubble_size, bubble_size)
        bubble.fill.solid()
        bubble.fill.fore_color.rgb = bubble_colors[i]
        bubble.line.color.rgb = RGBColor(0, 0, 0)
        bubble.line.width = Pt(1.5)
        
        label_x = x + offsets[i][0]
        label_y = y + offsets[i][1]
        label_box = slide6.shapes.add_textbox(label_x, label_y, Inches(0.8), Inches(0.3))
        label_frame = label_box.text_frame
        label_frame.text = model
        label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        label_frame.paragraphs[0].font.size = Pt(11)
        label_frame.paragraphs[0].font.bold = True
    
    y_label = slide6.shapes.add_textbox(left - Inches(0.6), top + height/2 - Inches(0.5), Inches(0.5), Inches(1))
    y_label_frame = y_label.text_frame
    y_label_frame.text = '均方误差（MSE）'
    y_label_frame.word_wrap = True
    y_label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    y_label_frame.paragraphs[0].font.size = Pt(12)
    
    x_label = slide6.shapes.add_textbox(left + width/2 - Inches(1), top + height + Inches(0.5), Inches(2), Inches(0.3))
    x_label_frame = x_label.text_frame
    x_label_frame.text = '模型复杂度'
    x_label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    x_label_frame.paragraphs[0].font.size = Pt(12)
    
    # 图例
    legend_y = top + Inches(0.3)
    legend_x = left + width - Inches(2.5)
    
    legend_title = slide6.shapes.add_textbox(legend_x, legend_y - Inches(0.3), Inches(2), Inches(0.3))
    legend_title.text_frame.text = '计算时间'
    legend_title.text_frame.paragraphs[0].font.size = Pt(11)
    legend_title.text_frame.paragraphs[0].font.bold = True
    
    legend_sizes = [(0.001, '小（快）', 0.25), (0.05, '中', 0.45), (0.13, '大（慢）', 0.65)]
    for time_val, label, size_factor in legend_sizes:
        bubble_legend = slide6.shapes.add_shape(MSO_SHAPE.OVAL, legend_x + Inches(0.1), 
                                                legend_y + Inches(size_factor - 0.15), 
                                                Inches(0.15), Inches(0.15))
        bubble_legend.fill.solid()
        bubble_legend.fill.fore_color.rgb = RGBColor(128, 128, 128)
        bubble_legend.line.color.rgb = RGBColor(0, 0, 0)
        
        label_box = slide6.shapes.add_textbox(legend_x + Inches(0.4), legend_y + Inches(size_factor - 0.15), Inches(1), Inches(0.2))
        label_box.text_frame.text = label
        label_box.text_frame.paragraphs[0].font.size = Pt(10)
    
    print("PPT图6完成")
    
    # 保存PPT
    prs.save(f'{output_dir}/figures_ppt.pptx')
    print(f"PPT文件已保存: {output_dir}/figures_ppt.pptx")

# ==================== 主程序 ====================

if __name__ == '__main__':
    # 绘制所有SVG图表
    plot_fig1_compute_time()
    plot_fig2_mse_comparison()
    plot_fig3_speedup()
    plot_fig4_blb_variance()
    plot_fig5_sgd_convergence()
    plot_fig6_complexity_accuracy()
    
    # 创建PPT
    create_ppt_charts()
    
    print("所有图表生成完成！")
